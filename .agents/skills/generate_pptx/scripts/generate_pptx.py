import os
import sys
import json
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 18 Color Palettes from pptx-official ──────────────────────────────────────
PALETTES = {
    "classic_blue": {"bg": "#F4F6F6", "text": "#1C2833", "primary": "#1C2833", "secondary": "#2E4053"},
    "teal_coral": {"bg": "#FFFFFF", "text": "#277884", "primary": "#277884", "secondary": "#FE4447"},
    "bold_red": {"bg": "#FFFFFF", "text": "#C0392B", "primary": "#C0392B", "secondary": "#F39C12"},
    "warm_blush": {"bg": "#FAF7F2", "text": "#A49393", "primary": "#A49393", "secondary": "#E8B4B8"},
    "burgundy_luxury": {"bg": "#FAF7F2", "text": "#5D1D2E", "primary": "#5D1D2E", "secondary": "#997929"},
    "deep_purple_emerald": {"bg": "#181B24", "text": "#FFFFFF", "primary": "#B165FB", "secondary": "#40695B"},
    "cream_forest_green": {"bg": "#FFE1C7", "text": "#40695B", "primary": "#40695B", "secondary": "#FCFCFC"},
    "pink_purple": {"bg": "#3D2F68", "text": "#FFFFFF", "primary": "#F8275B", "secondary": "#FF737D"},
    "lime_plum": {"bg": "#98ACB5", "text": "#7C3A5F", "primary": "#7C3A5F", "secondary": "#C5DE82"},
    "black_gold": {"bg": "#F4F6F6", "text": "#000000", "primary": "#000000", "secondary": "#BF9A4A"},
    "sage_terracotta": {"bg": "#F4F1DE", "text": "#2C2C2C", "primary": "#2C2C2C", "secondary": "#E07A5F"},
    "charcoal_red": {"bg": "#CCCBCB", "text": "#292929", "primary": "#292929", "secondary": "#E33737"},
    "vibrant_orange": {"bg": "#F2F2F2", "text": "#222831", "primary": "#222831", "secondary": "#F96D00"},
    "forest_green": {"bg": "#FFFFFF", "text": "#1E5128", "primary": "#1E5128", "secondary": "#4E9F3D"},
    "retro_rainbow": {"bg": "#FFFFFF", "text": "#722880", "primary": "#722880", "secondary": "#DEB600"},
    "vintage_earthy": {"bg": "#F4F1DE", "text": "#3A6B35", "primary": "#3A6B35", "secondary": "#E3B448"},
    "coastal_rose": {"bg": "#F3ECDC", "text": "#AD7670", "primary": "#AD7670", "secondary": "#B49886"},
    "orange_turquoise": {"bg": "#FCFCFC", "text": "#667C6F", "primary": "#667C6F", "secondary": "#FC993E"}
}

def parse_args():
    parser = argparse.ArgumentParser(description="Create custom PPTX presentation from JSON schema")
    parser.add_argument("--schema", required=True, help="Path to JSON schema file")
    parser.add_argument("--output", default="output.pptx", help="Path to output PPTX file")
    return parser.parse_args()

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    if len(hex_str) == 3:
        hex_str = ''.join(c*2 for c in hex_str)
    return RGBColor(*(int(hex_str[i:i+2], 16) for i in (0, 2, 4)))

def create_presentation(schema_path, output_path):
    with open(schema_path, 'r', encoding='utf-8') as f:
        schema = json.load(f)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    theme = schema.get('theme', {})
    palette_name = theme.get('palette')
    
    # Resolve color presets
    if palette_name and palette_name in PALETTES:
        preset = PALETTES[palette_name]
        bg_hex = preset["bg"]
        text_hex = preset["text"]
        primary_hex = preset["primary"]
        secondary_hex = preset["secondary"]
    else:
        bg_hex = theme.get('backgroundColor', '#ffffff')
        text_hex = theme.get('textColor', '#1e293b')
        primary_hex = theme.get('primaryColor', '#6366f1')
        secondary_hex = theme.get('secondaryColor', '#f43f5e')
    
    bg_color = hex_to_rgb(bg_hex)
    text_color = hex_to_rgb(text_hex)
    primary_color = hex_to_rgb(primary_hex)
    secondary_color = hex_to_rgb(secondary_hex)
    
    font_name = theme.get('fontName', 'Calibri')
    
    blank_slide_layout = prs.slide_layouts[6]
    
    for slide_data in schema.get('slides', []):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Slide overrides or fallbacks
        slide_bg_hex = slide_data.get('backgroundColor', bg_hex)
        slide_bg_color = hex_to_rgb(slide_bg_hex)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = slide_bg_color
        
        slide_text_hex = slide_data.get('textColor', text_hex)
        slide_text_color = hex_to_rgb(slide_text_hex)
        
        slide_type = slide_data.get('type', 'content')
        
        if slide_type == 'title':
            # Title slide layout
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = slide_data.get('title', '')
            p.font.name = font_name
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = slide_text_color
            
            sub = slide_data.get('subtitle', '')
            if sub:
                sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1.5))
                tf_sub = sub_box.text_frame
                tf_sub.word_wrap = True
                p_sub = tf_sub.paragraphs[0]
                p_sub.alignment = PP_ALIGN.CENTER
                p_sub.text = sub
                p_sub.font.name = font_name
                p_sub.font.size = Pt(20)
                p_sub.font.color.rgb = secondary_color
                
        elif slide_type == 'header':
            # Section breaker slide
            # Highlight with a strong primary background
            slide_header_bg = hex_to_rgb(slide_data.get('backgroundColor', primary_hex))
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = slide_header_bg
            
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(2.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = slide_data.get('title', '')
            p.font.name = font_name
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb('#ffffff')
            
        elif slide_type == 'two_columns':
            # Left aligned slide title
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(0.9))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get('title', '')
            p.font.name = font_name
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = primary_color
            
            # Draw optional accent line under title (design preference)
            if slide_data.get('underlineAccent', True):
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(2.0), Inches(0.04))
                line.fill.solid()
                line.fill.fore_color.rgb = secondary_color
                line.line.fill.background()
            
            # Layout splitting
            split = slide_data.get('split', '50_50')
            total_width = 11.733
            gap = 0.633
            
            if split == '30_70':
                left_w = total_width * 0.3 - gap/2
                right_w = total_width * 0.7 - gap/2
            elif split == '40_60':
                left_w = total_width * 0.4 - gap/2
                right_w = total_width * 0.6 - gap/2
            elif split == '70_30':
                left_w = total_width * 0.7 - gap/2
                right_w = total_width * 0.3 - gap/2
            elif split == '60_40':
                left_w = total_width * 0.6 - gap/2
                right_w = total_width * 0.4 - gap/2
            else: # '50_50'
                left_w = total_width * 0.5 - gap/2
                right_w = total_width * 0.5 - gap/2
                
            left_x = Inches(0.8)
            right_x = Inches(0.8) + Inches(left_w + gap)
            
            # Left column
            left_col = slide_data.get('leftColumn', {})
            left_box = slide.shapes.add_textbox(left_x, Inches(1.8), Inches(left_w), Inches(4.8))
            tf_left = left_box.text_frame
            tf_left.word_wrap = True
            
            col_title = left_col.get('title', '')
            first_p = True
            if col_title:
                p_lt = tf_left.paragraphs[0]
                p_lt.text = col_title
                p_lt.font.name = font_name
                p_lt.font.size = Pt(20)
                p_lt.font.bold = True
                p_lt.font.color.rgb = slide_text_color
                first_p = False
                
            for bullet in left_col.get('bullets', []):
                p_b = tf_left.add_paragraph() if not first_p else tf_left.paragraphs[0]
                first_p = False
                p_b.text = "• " + bullet
                p_b.font.name = font_name
                p_b.font.size = Pt(15)
                p_b.font.color.rgb = slide_text_color
                p_b.space_before = Pt(6)
                
            # Right column
            right_col = slide_data.get('rightColumn', {})
            right_box = slide.shapes.add_textbox(right_x, Inches(1.8), Inches(right_w), Inches(4.8))
            tf_right = right_box.text_frame
            tf_right.word_wrap = True
            
            col_title = right_col.get('title', '')
            first_p = True
            if col_title:
                p_rt = tf_right.paragraphs[0]
                p_rt.text = col_title
                p_rt.font.name = font_name
                p_rt.font.size = Pt(20)
                p_rt.font.bold = True
                p_rt.font.color.rgb = slide_text_color
                first_p = False
                
            for bullet in right_col.get('bullets', []):
                p_b = tf_right.add_paragraph() if not first_p else tf_right.paragraphs[0]
                first_p = False
                p_b.text = "• " + bullet
                p_b.font.name = font_name
                p_b.font.size = Pt(15)
                p_b.font.color.rgb = slide_text_color
                p_b.space_before = Pt(6)
                
        else: # 'content'
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(0.9))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get('title', '')
            p.font.name = font_name
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = primary_color
            
            # Draw optional accent line under title
            if slide_data.get('underlineAccent', True):
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(2.0), Inches(0.04))
                line.fill.solid()
                line.fill.fore_color.rgb = secondary_color
                line.line.fill.background()
            
            # Content Box
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            
            bullets = slide_data.get('bullets', [])
            first_p = True
            for bullet in bullets:
                p_b = tf_content.add_paragraph() if not first_p else tf_content.paragraphs[0]
                first_p = False
                p_b.text = "• " + bullet
                p_b.font.name = font_name
                p_b.font.size = Pt(16)
                p_b.font.color.rgb = slide_text_color
                p_b.space_before = Pt(10)
                
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    args = parse_args()
    create_presentation(args.schema, args.output)
